# -*- coding: utf-8 -*-
"""Shared helpers for building ToxAgent decks on Slides_template.pptx."""
import copy
import os
import tempfile
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ---------------------------------------------------------------- palette
NAVY   = RGBColor(0x00, 0x20, 0x60)   # slide titles
DARK   = RGBColor(0x0E, 0x28, 0x41)   # theme dark-2, bars / section bg
BLUE   = RGBColor(0x15, 0x60, 0x82)   # accent1
ORANGE = RGBColor(0xE9, 0x71, 0x32)   # accent2
GREEN  = RGBColor(0x19, 0x6B, 0x24)   # accent3
CYAN   = RGBColor(0x0F, 0x9E, 0xD5)   # accent4
RED    = RGBColor(0xC0, 0x00, 0x00)
GRAY   = RGBColor(0x59, 0x59, 0x59)
LGRAY  = RGBColor(0xF2, 0xF2, 0xF2)
MGRAY  = RGBColor(0xD9, 0xD9, 0xD9)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
PALE   = RGBColor(0x9F, 0xC5, 0xE8)
PALEBG = RGBColor(0xE7, 0xEF, 0xF7)
PALEOR = RGBColor(0xFD, 0xEC, 0xE0)
PALEGR = RGBColor(0xE4, 0xF0, 0xE5)
BLACK  = RGBColor(0x00, 0x00, 0x00)

FONT = "Times New Roman"
MONO = "Consolas"

TEMPLATE = "Slides_template.pptx"
BLANK_LAYOUT = 6


# ---------------------------------------------------------------- chrome
class Deck:
    def __init__(self, template=TEMPLATE, left_label="ToxAgent Harness",
                 right_label="NEU Bio-Research Team"):
        self.prs = Presentation(template)
        self.left_label = left_label
        self.right_label = right_label
        # keep copies of the template's chrome before deleting sample slides
        s_title = self.prs.slides[0]
        s_body = self.prs.slides[2]   # 'About' slide: chrome only, no agenda list
        # pictures cannot be carried by raw XML copy (their r:embed rels would
        # dangle once the source slide part is dropped) — re-add them by blob.
        self._picdir = tempfile.mkdtemp(prefix="deck_pics_")
        self._title_pics = []
        for sh in s_title.shapes:
            if sh.shape_type == 13:
                fn = os.path.join(self._picdir,
                                  "p%d.%s" % (len(self._title_pics), sh.image.ext))
                with open(fn, "wb") as fh:
                    fh.write(sh.image.blob)
                self._title_pics.append((fn, sh.left, sh.top, sh.width, sh.height))
        self._title_xml = [copy.deepcopy(sh._element)
                           for sh in s_title.shapes if sh.shape_type != 13]
        self._chrome_xml = [copy.deepcopy(sh._element)
                            for sh in s_body.shapes if sh.shape_type != 13]
        self._drop_all_slides()

    def _drop_all_slides(self):
        xml_slides = self.prs.slides._sldIdLst
        for sld in list(xml_slides):
            rId = sld.get(qn('r:id'))
            self.prs.part.drop_rel(rId)
            xml_slides.remove(sld)

    # -- slide factory ----------------------------------------------------
    def blank(self):
        return self.prs.slides.add_slide(self.prs.slide_layouts[BLANK_LAYOUT])

    def content(self, title=None, title_size=32):
        """Content slide with the template footer bar + optional big title."""
        s = self.blank()
        tree = s.shapes._spTree
        for el in self._chrome_xml:
            tree.append(copy.deepcopy(el))
        # rewrite the two chrome labels
        for sh in s.shapes:
            if sh.has_text_frame and sh.text_frame.text.strip():
                t = sh.text_frame.text.strip()
                if t in ("Trong-Nghia Nguyen", "NEU Bio-Research Team"):
                    _set_run(sh, self.right_label)
                elif t in ("Time-series", "Intro to AI", "ToxAgent Harness"):
                    _set_run(sh, self.left_label)
                elif not sh.is_placeholder:   # template's own title/body boxes
                    sh._element.getparent().remove(sh._element)
        if title is not None:
            tb = s.shapes.add_textbox(Inches(0.30), Inches(0.17), Inches(12.73), Inches(0.72))
            tf = tb.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            r = p.add_run(); r.text = title
            r.font.name = FONT; r.font.size = Pt(title_size)
            r.font.bold = True; r.font.color.rgb = NAVY
        return s

    def title_slide(self, title, subtitle_lines, badge, corner_lines):
        s = self.blank()
        tree = s.shapes._spTree
        for el in self._title_xml:
            tree.append(copy.deepcopy(el))
        shapes = list(s.shapes)
        for sh in shapes:
            if not sh.has_text_frame:
                continue
            t = sh.text_frame.text.strip()
            if t.startswith("Topic Name") or "Đề xuất kiến trúc" in t:
                _set_para_lines(sh, [(title, 30, True, DARK)], align=PP_ALIGN.CENTER)
            elif t.startswith("Group 1") or "NEU Bio-Research" in t:
                _set_para_lines(sh, subtitle_lines, align=PP_ALIGN.CENTER)
            elif t.startswith("Progress Report") or "Architecture" in t:
                _set_para_lines(sh, [(badge, 20, True, BLACK)], align=PP_ALIGN.CENTER)
            elif t.startswith("Time-series,") or t.startswith("ToxAgent,"):
                _set_para_lines(sh, [(l, 14, False, WHITE) for l in corner_lines],
                                align=PP_ALIGN.CENTER, font="Arial")
        for fn, l, t, w, h in self._title_pics:
            s.shapes.add_picture(fn, l, t, w, h)
        return s

    def section(self, kicker, title, subtitle=""):
        s = self.content(title=None)
        band = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(2.35),
                                  Inches(13.333), Inches(1.90))
        band.fill.solid(); band.fill.fore_color.rgb = DARK
        band.line.fill.background()
        band.shadow.inherit = False
        band.text_frame.text = ""
        textbox(s, 0.90, 2.50, 11.55, 0.55, [(kicker, 20, True, PALE)], align=PP_ALIGN.CENTER)
        textbox(s, 0.60, 2.98, 12.13, 1.00, [(title, 40, True, WHITE)], align=PP_ALIGN.CENTER)
        if subtitle:
            textbox(s, 1.10, 4.50, 11.13, 1.10, [(subtitle, 22, False, GRAY)],
                    align=PP_ALIGN.CENTER)
        return s

    def save(self, path):
        self.prs.save(path)
        return path


# ---------------------------------------------------------------- text
def _set_run(shape, text):
    tf = shape.text_frame
    p = tf.paragraphs[0]
    if p.runs:
        p.runs[0].text = text
        for extra in p.runs[1:]:
            extra._r.getparent().remove(extra._r)
    else:
        r = p.add_run(); r.text = text


def _set_para_lines(shape, lines, align=PP_ALIGN.LEFT, font=FONT):
    """lines: list of (text, size, bold, color) — replaces the shape's text."""
    tf = shape.text_frame
    tf.word_wrap = True
    for p in list(tf.paragraphs)[1:]:
        p._p.getparent().remove(p._p)
    p0 = tf.paragraphs[0]
    for r in list(p0.runs):
        r._r.getparent().remove(r._r)
    for i, item in enumerate(lines):
        text, size, bold, color = item
        p = p0 if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run(); r.text = text
        r.font.name = font; r.font.size = Pt(size)
        r.font.bold = bold
        if color is not None:
            r.font.color.rgb = color


def textbox(slide, x, y, w, h, lines, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
            font=FONT, space_after=4, line_spacing=None):
    """lines: list of (text, size, bold, color) or (text, size, bold, color, indent)."""
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Inches(0.04)
    tf.margin_top = tf.margin_bottom = 0
    for i, item in enumerate(lines):
        text, size, bold, color = item[:4]
        indent = item[4] if len(item) > 4 else 0
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space_after)
        if line_spacing:
            p.line_spacing = line_spacing
        if indent:
            p.level = indent
        r = p.add_run(); r.text = text
        r.font.name = font; r.font.size = Pt(size); r.font.bold = bold
        if color is not None:
            r.font.color.rgb = color
    return tb


def bullets(slide, x, y, w, h, items, size=21, color=DARK, bullet="\u2022 ",
            gap=9, align=PP_ALIGN.LEFT):
    """items: str, or (str, bold), or (str, bold, color), or ('', ...) for spacer."""
    lines = []
    for it in items:
        if isinstance(it, str):
            txt, bold, col = it, False, color
        elif len(it) == 2:
            txt, bold = it; col = color
        else:
            txt, bold, col = it
        if txt == "":
            lines.append((" ", 8, False, None))
        else:
            pre = bullet if not txt.startswith(("\u2022", "\u2013", "  ")) else ""
            lines.append((pre + txt, size, bold, col))
    return textbox(slide, x, y, w, h, lines, align=align, space_after=gap)


# ---------------------------------------------------------------- shapes
def box(slide, x, y, w, h, lines, fill=PALEBG, line=None, radius=True,
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, line_w=1.25):
    shp = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(h))
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line; shp.line.width = Pt(line_w)
    shp.shadow.inherit = False
    if radius:
        try:
            shp.adjustments[0] = 0.09
        except Exception:
            pass
    tf = shp.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Inches(0.08)
    tf.margin_top = tf.margin_bottom = Inches(0.04)
    if isinstance(lines, str):
        lines = [(lines, 18, True, DARK)]
    for i, item in enumerate(lines):
        text, size, bold, color = item[:4]
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(2)
        r = p.add_run(); r.text = text
        r.font.name = FONT; r.font.size = Pt(size); r.font.bold = bold
        if color is not None:
            r.font.color.rgb = color
    return shp


def stat_card(slide, x, y, w, h, big, small_lines, big_color=RED,
              fill=LGRAY, border=MGRAY):
    box(slide, x, y, w, h, [], fill=fill, line=border)
    textbox(slide, x + 0.10, y + 0.12, w - 0.20, 0.70,
            [(big, 40, True, big_color)], align=PP_ALIGN.CENTER)
    textbox(slide, x + 0.10, y + 0.86, w - 0.20, h - 0.95,
            [(s, 15, False, GRAY) for s in small_lines],
            align=PP_ALIGN.CENTER, space_after=1)


def banner(slide, x, y, w, h, text, size=19, fill=DARK, color=WHITE, bold=True):
    return box(slide, x, y, w, h, [(text, size, bold, color)], fill=fill)


def arrow(slide, x, y, w, h, color=BLUE, direction="right"):
    shape = {"right": MSO_SHAPE.RIGHT_ARROW, "down": MSO_SHAPE.DOWN_ARROW,
             "left": MSO_SHAPE.LEFT_ARROW, "up": MSO_SHAPE.UP_ARROW}[direction]
    a = slide.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    a.fill.solid(); a.fill.fore_color.rgb = color
    a.line.fill.background()
    a.shadow.inherit = False
    return a


def chevron(slide, x, y, w, h, lines, fill=BLUE):
    shp = slide.shapes.add_shape(MSO_SHAPE.CHEVRON, Inches(x), Inches(y),
                                 Inches(w), Inches(h))
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    shp.line.fill.background(); shp.shadow.inherit = False
    tf = shp.text_frame; tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    for i, (text, size, bold, color) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = text
        r.font.name = FONT; r.font.size = Pt(size); r.font.bold = bold
        r.font.color.rgb = color
    return shp


def picture(slide, path, x, y, w=None, h=None, border=True):
    kw = {}
    if w: kw["width"] = Inches(w)
    if h: kw["height"] = Inches(h)
    pic = slide.shapes.add_picture(path, Inches(x), Inches(y), **kw)
    if border:
        pic.line.color.rgb = MGRAY
        pic.line.width = Pt(1)
    return pic


# ---------------------------------------------------------------- table
def table(slide, x, y, w, h, data, col_w=None, font_size=16,
          header_fill=DARK, header_color=WHITE, body_color=DARK,
          zebra=True, bold_first_col=False, row_h=None, align_first=PP_ALIGN.LEFT):
    rows, cols = len(data), len(data[0])
    gf = slide.shapes.add_table(rows, cols, Inches(x), Inches(y),
                                Inches(w), Inches(h))
    tbl = gf.table
    tbl.first_row = True
    # kill the banded default style so our fills show
    tbl._tbl.tblPr.set('firstRow', '1')
    tbl._tbl.tblPr.set('bandRow', '0')
    if col_w:
        total = sum(col_w)
        for i, cw in enumerate(col_w):
            tbl.columns[i].width = Emu(int(Inches(w) * cw / total))
    if row_h:
        for i, rh in enumerate(row_h):
            tbl.rows[i].height = Inches(rh)
    for ri, row in enumerate(data):
        for ci, val in enumerate(row):
            cell = tbl.cell(ri, ci)
            cell.margin_left = Inches(0.09)
            cell.margin_right = Inches(0.07)
            cell.margin_top = Inches(0.045)
            cell.margin_bottom = Inches(0.045)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            if ri == 0:
                cell.fill.solid(); cell.fill.fore_color.rgb = header_fill
            elif zebra and ri % 2 == 0:
                cell.fill.solid(); cell.fill.fore_color.rgb = LGRAY
            else:
                cell.fill.solid(); cell.fill.fore_color.rgb = WHITE
            tf = cell.text_frame
            tf.word_wrap = True
            parts = val if isinstance(val, list) else [val]
            for pi, ptxt in enumerate(parts):
                p = tf.paragraphs[0] if pi == 0 else tf.add_paragraph()
                p.alignment = align_first if ci == 0 else PP_ALIGN.LEFT
                p.space_after = Pt(0)
                bold = (ri == 0) or (bold_first_col and ci == 0)
                colr = header_color if ri == 0 else body_color
                txt = ptxt
                if isinstance(ptxt, tuple):
                    txt, extra = ptxt
                    if extra == "b": bold = True
                    elif isinstance(extra, RGBColor): colr = extra
                r = p.add_run(); r.text = txt
                r.font.name = FONT
                r.font.size = Pt(font_size if ri else font_size)
                r.font.bold = bold
                r.font.color.rgb = colr
    return gf
